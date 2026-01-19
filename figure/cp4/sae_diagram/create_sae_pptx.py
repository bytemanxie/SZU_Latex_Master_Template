#!/usr/bin/env python3
"""
创建完整 SAE 模块（CoordSaeLayer）结构图 PPT
残差用分叉箭头 + 号表示
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(80, 80, 80)

CX = 5.0
BW = 1.7
BH = 0.35
SBW = 0.75
SBH = 0.3
GAP = 0.12
RES_X = CX - 1.4  # 残差线 X 位置


def box(left, top, w, h, text, size=10):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BLACK
    s.line.width = Pt(1.5)
    p = s.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size)
    p.font.name = "Times New Roman"
    p.font.color.rgb = BLACK
    s.text_frame_anchor = MSO_ANCHOR.MIDDLE
    return s


def circle(left, top, sz, text):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(sz), Inches(sz))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BLACK
    s.line.width = Pt(1.5)
    p = s.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLACK
    s.text_frame_anchor = MSO_ANCHOR.MIDDLE
    return s


def label(left, top, text, size=9, italic=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(1.5), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.italic = italic
    p.font.name = "Times New Roman"
    p.font.color.rgb = GRAY
    return tb


def line(x1, y1, x2, y2, w=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK
    c.line.width = Pt(w)
    return c


# ============ 绘制 SAE 模块 ============
y = 0.3
cs = 0.25  # 圆圈大小

# Input
label(CX - 0.15, y, "X", size=12, italic=False)
y += 0.28
line(CX, y, CX, y + GAP)
y += GAP

# === Stage 1: CoordAtt ===
# 分叉点（残差1开始）
res1_start = y
line(RES_X, y, CX, y)  # 水平分叉线
line(RES_X, y, RES_X, y)  # 残差起点
line(CX, y, CX, y + GAP)
y += GAP

# X Avg Pool | Y Avg Pool
pg = 0.25
pl = CX - pg/2 - SBW
pr = CX + pg/2
line(CX, y - 0.03, CX, y)
line(pl + SBW/2, y - 0.03, pr + SBW/2, y - 0.03)
line(pl + SBW/2, y - 0.03, pl + SBW/2, y)
line(pr + SBW/2, y - 0.03, pr + SBW/2, y)
box(pl, y, SBW, SBH, "X Avg Pool", size=8)
box(pr, y, SBW, SBH, "Y Avg Pool", size=8)
label(pl - 0.75, y + 0.02, "C×H×1", size=8)
label(pr + SBW + 0.05, y + 0.02, "C×1×W", size=8)
y += SBH + GAP

# 汇合
line(pl + SBW/2, y - GAP, pl + SBW/2, y - 0.03)
line(pr + SBW/2, y - GAP, pr + SBW/2, y - 0.03)
line(pl + SBW/2, y - 0.03, pr + SBW/2, y - 0.03)
line(CX, y - 0.03, CX, y)

# Concat + Conv2d
box(CX - BW/2, y, BW, BH, "Concat + Conv2d")
y += BH + GAP
line(CX, y - GAP, CX, y)

# BN + Non-linear
box(CX - BW/2, y, BW, BH, "BatchNorm + Non-linear")
y += BH + GAP
line(CX, y - GAP, CX, y - 0.03)

# Split
label(CX - 0.15, y - 0.1, "split", size=8)
line(pl + SBW/2, y - 0.03, pr + SBW/2, y - 0.03)
line(pl + SBW/2, y - 0.03, pl + SBW/2, y)
line(pr + SBW/2, y - 0.03, pr + SBW/2, y)

# Conv2d | Conv2d
box(pl, y, SBW, SBH, "Conv2d", size=9)
box(pr, y, SBW, SBH, "Conv2d", size=9)
y += SBH + GAP

# Sigmoid | Sigmoid
line(pl + SBW/2, y - GAP, pl + SBW/2, y)
line(pr + SBW/2, y - GAP, pr + SBW/2, y)
box(pl, y, SBW, SBH, "Sigmoid", size=9)
box(pr, y, SBW, SBH, "Sigmoid", size=9)
y += SBH + GAP

# 汇合到 Re-weight
line(pl + SBW/2, y - GAP, pl + SBW/2, y + 0.05)
line(pr + SBW/2, y - GAP, pr + SBW/2, y + 0.05)
line(pl + SBW/2, y + 0.05, CX - BW/2, y + 0.05)
line(pr + SBW/2, y + 0.05, CX + BW/2, y + 0.05)
y += 0.05

# Re-weight
box(CX - BW/2, y, BW, BH, "Re-weight")
y += BH + GAP
line(CX, y - GAP, CX, y)

# (+) 残差1加法
res1_end = y + cs/2
line(RES_X, res1_start, RES_X, res1_end)  # 残差垂直线
line(RES_X, res1_end, CX - cs/2, res1_end)  # 残差水平线到+号
circle(CX - cs/2, y, cs, "+")
y += cs + GAP
line(CX, y, CX, y + GAP)
y += GAP

# === Stage 2: Conv Block ===
# 分叉点（残差2开始）
res2_start = y
line(RES_X, y, CX, y)
line(CX, y, CX, y + GAP)
y += GAP

# Conv Block
box(CX - BW/2, y, BW, BH * 1.5, "3×3 Conv+BN+ReLU ×2", size=9)
y += BH * 1.5 + GAP
line(CX, y - GAP, CX, y)

# (+) 残差2加法
res2_end = y + cs/2
line(RES_X, res2_start, RES_X, res2_end)
line(RES_X, res2_end, CX - cs/2, res2_end)
circle(CX - cs/2, y, cs, "+")
y += cs + GAP
line(CX, y, CX, y + GAP)
y += GAP

# === Stage 3: SaELayer ===
# 分叉点（Scale输入）
scale_input_start = y
scale_x = RES_X
line(scale_x, y, CX, y)
line(CX, y, CX, y + GAP)
y += GAP

# Global Avg Pool
box(CX - BW/2, y, BW, BH, "Global Avg Pool")
label(CX + BW/2 + 0.1, y + 0.03, "1×1×C")
y += BH + GAP
line(CX, y - GAP, CX, y - 0.03)

# 4个并行 FC+ReLU
fc_gap = 0.08
fc_w = 0.6
fc_total = 4 * fc_w + 3 * fc_gap
fc_start = CX - fc_total / 2

line(fc_start + fc_w/2, y - 0.03, fc_start + fc_total - fc_w/2, y - 0.03)
for i in range(4):
    fx = fc_start + i * (fc_w + fc_gap)
    line(fx + fc_w/2, y - 0.03, fx + fc_w/2, y)
    box(fx, y, fc_w, SBH, f"FC+ReLU", size=7)
    line(fx + fc_w/2, y + SBH, fx + fc_w/2, y + SBH + 0.08)

label(fc_start + fc_total + 0.05, y + 0.02, "1×1×C/r ×4", size=8)
y += SBH + 0.08

# 汇合
line(fc_start + fc_w/2, y, fc_start + fc_total - fc_w/2, y)
line(CX, y, CX, y + 0.05)
y += 0.05

# Concat
box(CX - 0.5, y, 1.0, SBH, "Concat", size=9)
y += SBH + GAP
line(CX, y - GAP, CX, y)

# FC + Sigmoid
box(CX - BW/2, y, BW, BH, "FC + Sigmoid")
label(CX + BW/2 + 0.1, y + 0.03, "1×1×C")
y += BH + GAP
line(CX, y - GAP, CX, y)

# Scale (×) - 与输入相乘
scale_end = y + cs/2
line(scale_x, scale_input_start, scale_x, scale_end)
line(scale_x, scale_end, CX - cs/2, scale_end)
circle(CX - cs/2, y, cs, "×")
label(CX + cs/2 + 0.08, y + 0.02, "Scale", size=10, italic=False)
y += cs + GAP
line(CX, y, CX, y + GAP)
y += GAP

# Output
label(CX - 0.15, y, "X̃", size=12, italic=False)
y += 0.3

# 模块名称
label(CX - 0.85, y, "SAE Module", size=12, italic=False)

# 保存
output_path = '/Users/xiezhijie/Documents/SZU_Latex_Master_Template/figure/cp4/sae_diagram/fig4-1_sae_module.pptx'
prs.save(output_path)
print(f"已保存: {output_path}")
