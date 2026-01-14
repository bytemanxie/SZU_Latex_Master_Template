#!/usr/bin/env python3
"""
生成双层路由注意力机制示意图 PPT
图 3-4: 双层路由注意力机制示意图
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 颜色定义
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_QUERY = RGBColor(255, 107, 107)  # 红色 - 查询窗口
COLOR_QUERY_BORDER = RGBColor(192, 0, 0)
COLOR_SELECTED = RGBColor(169, 209, 142)  # 绿色 - 被选中窗口
COLOR_SELECTED_BORDER = RGBColor(84, 130, 53)
COLOR_ATTENTION_BG = RGBColor(222, 235, 247)  # 浅蓝色
COLOR_ATTENTION_BORDER = RGBColor(46, 117, 182)
COLOR_ARROW = RGBColor(89, 89, 89)  # 深灰色
COLOR_LINE = RGBColor(84, 130, 53)  # 深绿色连接线


def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    # 设置幻灯片大小为宽屏
    prs.slide_width = Cm(28)
    prs.slide_height = Cm(16)
    return prs


def add_rectangle(slide, left, top, width, height, fill_color, line_color, line_width=Pt(1)):
    """添加矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape


def add_rounded_rectangle(slide, left, top, width, height, fill_color, line_color, line_width=Pt(1)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=Pt(12), bold=False, alignment=PP_ALIGN.CENTER):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.name = "宋体"
    p.alignment = alignment
    return txBox


def add_arrow(slide, left, top, width, height):
    """添加箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ARROW
    shape.line.color.rgb = COLOR_ARROW
    return shape


def draw_grid(slide, start_left, start_top, cell_size, rows, cols, query_pos=None, selected_positions=None):
    """绘制窗口网格
    
    Args:
        query_pos: 查询窗口位置 (row, col)，从0开始
        selected_positions: 被选中窗口位置列表 [(row, col), ...]
    """
    shapes = []
    for row in range(rows):
        for col in range(cols):
            left = start_left + col * cell_size
            top = start_top + row * cell_size
            
            # 确定颜色
            if query_pos and (row, col) == query_pos:
                fill_color = COLOR_QUERY
                line_color = COLOR_QUERY_BORDER
                line_width = Pt(2)
            elif selected_positions and (row, col) in selected_positions:
                fill_color = COLOR_SELECTED
                line_color = COLOR_SELECTED_BORDER
                line_width = Pt(1.5)
            else:
                fill_color = COLOR_WHITE
                line_color = COLOR_BLACK
                line_width = Pt(1)
            
            shape = add_rectangle(slide, left, top, cell_size, cell_size, 
                                 fill_color, line_color, line_width)
            shapes.append(shape)
    
    return shapes


def draw_connection_line(slide, start_shape, end_shape):
    """绘制连接线（使用直线连接器）"""
    # 获取形状中心点
    start_x = start_shape.left + start_shape.width // 2
    start_y = start_shape.top + start_shape.height // 2
    end_x = end_shape.left + end_shape.width // 2
    end_y = end_shape.top + end_shape.height // 2
    
    # 添加直线
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = COLOR_LINE
    connector.line.width = Pt(1.5)
    return connector


def create_diagram_slide(prs):
    """创建双层路由注意力机制示意图幻灯片"""
    
    # 添加空白幻灯片
    blank_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(blank_layout)
    
    # ========== 参数设置 ==========
    cell_size = Cm(0.9)
    grid_rows = 4
    grid_cols = 4
    
    # 三个部分的起始位置
    section1_left = Cm(2)    # 窗口划分
    section2_left = Cm(10)   # 稀疏选择
    section3_left = Cm(18)   # 局部注意力
    
    grid_top = Cm(4)
    
    # ========== 第一部分：窗口划分 ==========
    # 绘制网格，第2个位置(0,1)为查询窗口
    grid1_shapes = draw_grid(slide, section1_left, grid_top, cell_size, 
                             grid_rows, grid_cols, query_pos=(0, 1))
    
    # 添加标题
    add_text_box(slide, section1_left - Cm(0.5), grid_top + cell_size * grid_rows + Cm(0.3),
                 Cm(5), Cm(1), "窗口划分", Pt(14), bold=True)
    
    # ========== 第一个箭头 ==========
    arrow1_left = section1_left + cell_size * grid_cols + Cm(0.5)
    add_arrow(slide, arrow1_left, grid_top + cell_size * 1.5, Cm(1.8), Cm(0.8))
    
    # ========== 第二部分：稀疏选择 ==========
    # 查询窗口(0,1)，被选中窗口(1,3), (2,1), (3,2)
    selected_pos = [(1, 3), (2, 1), (3, 2)]
    grid2_shapes = draw_grid(slide, section2_left, grid_top, cell_size,
                             grid_rows, grid_cols, query_pos=(0, 1), 
                             selected_positions=selected_pos)
    
    # 绘制连接线
    query_idx = 0 * grid_cols + 1  # 查询窗口在shapes中的索引
    query_shape = grid2_shapes[query_idx]
    
    for pos in selected_pos:
        target_idx = pos[0] * grid_cols + pos[1]
        target_shape = grid2_shapes[target_idx]
        draw_connection_line(slide, query_shape, target_shape)
    
    # 添加标题
    add_text_box(slide, section2_left - Cm(0.5), grid_top + cell_size * grid_rows + Cm(0.3),
                 Cm(5), Cm(0.6), "稀疏选择", Pt(14), bold=True)
    add_text_box(slide, section2_left - Cm(0.5), grid_top + cell_size * grid_rows + Cm(0.9),
                 Cm(5), Cm(0.6), "(TopK 路由)", Pt(11), bold=False)
    
    # ========== 第二个箭头 ==========
    arrow2_left = section2_left + cell_size * grid_cols + Cm(0.5)
    add_arrow(slide, arrow2_left, grid_top + cell_size * 1.5, Cm(1.8), Cm(0.8))
    
    # ========== 第三部分：局部注意力 ==========
    # 圆角矩形框
    attention_box = add_rounded_rectangle(
        slide, section3_left, grid_top - Cm(0.3), 
        Cm(4.5), Cm(4.5),
        COLOR_ATTENTION_BG, COLOR_ATTENTION_BORDER, Pt(1.5)
    )
    
    # 内部文字
    formula_texts = [
        ("Q × K", Cm(0.3)),
        ("↓", Cm(0.9)),
        ("Softmax", Cm(1.5)),
        ("↓", Cm(2.1)),
        ("A × V", Cm(2.7)),
        ("↓", Cm(3.3)),
        ("Output", Cm(3.9)),
    ]
    
    for text, offset in formula_texts:
        add_text_box(slide, section3_left + Cm(0.8), grid_top - Cm(0.3) + offset,
                     Cm(3), Cm(0.6), text, Pt(13), bold=False)
    
    # 添加标题
    add_text_box(slide, section3_left, grid_top + cell_size * grid_rows + Cm(0.3),
                 Cm(5), Cm(1), "局部注意力", Pt(14), bold=True)
    
    # ========== 添加图例 ==========
    legend_left = Cm(22)
    legend_top = Cm(10)
    legend_size = Cm(0.5)
    
    # 图例标题
    add_text_box(slide, legend_left, legend_top - Cm(0.6), Cm(4), Cm(0.5), 
                 "图例", Pt(11), bold=True, alignment=PP_ALIGN.LEFT)
    
    # 查询窗口图例
    add_rectangle(slide, legend_left, legend_top, legend_size, legend_size,
                  COLOR_QUERY, COLOR_QUERY_BORDER, Pt(1))
    add_text_box(slide, legend_left + Cm(0.7), legend_top - Cm(0.05), Cm(3), Cm(0.5),
                 "查询窗口", Pt(10), alignment=PP_ALIGN.LEFT)
    
    # 被选中窗口图例
    add_rectangle(slide, legend_left, legend_top + Cm(0.7), legend_size, legend_size,
                  COLOR_SELECTED, COLOR_SELECTED_BORDER, Pt(1))
    add_text_box(slide, legend_left + Cm(0.7), legend_top + Cm(0.65), Cm(3), Cm(0.5),
                 "被选中窗口", Pt(10), alignment=PP_ALIGN.LEFT)
    
    # 普通窗口图例
    add_rectangle(slide, legend_left, legend_top + Cm(1.4), legend_size, legend_size,
                  COLOR_WHITE, COLOR_BLACK, Pt(1))
    add_text_box(slide, legend_left + Cm(0.7), legend_top + Cm(1.35), Cm(3), Cm(0.5),
                 "忽略窗口", Pt(10), alignment=PP_ALIGN.LEFT)
    
    # ========== 添加两层标注 ==========
    # 虚线分隔说明 (用文本代替)
    add_text_box(slide, Cm(2), Cm(12), Cm(12), Cm(0.6),
                 "← 第一层：区域级路由（粗粒度）→", Pt(11))
    add_text_box(slide, Cm(16), Cm(12), Cm(8), Cm(0.6),
                 "← 第二层：细粒度注意力→", Pt(11))
    
    # ========== 添加图标题 ==========
    add_text_box(slide, Cm(0), Cm(14), Cm(28), Cm(1),
                 "图 3-4 双层路由注意力机制示意图", Pt(12), bold=True)
    
    return slide


def main():
    """主函数"""
    print("正在创建 PPT...")
    
    # 创建演示文稿
    prs = create_presentation()
    
    # 创建示意图幻灯片
    create_diagram_slide(prs)
    
    # 保存文件
    output_path = "/Users/xiezhijie/Documents/SZU_Latex_Master_Template/figure/cp3/fig3-4_biformer_diagram.pptx"
    prs.save(output_path)
    
    print(f"PPT 已保存到: {output_path}")
    print("请用 PowerPoint 打开并导出为 PNG 图片")


if __name__ == "__main__":
    main()
